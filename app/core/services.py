import importlib
import inspect
import logging
import os
import pkgutil
import uuid
from datetime import datetime
from pathlib import Path
from tempfile import NamedTemporaryFile
from time import perf_counter
from typing import Optional

import pypdfium2 as pdfium
from django.conf import settings
from pptx import Presentation

import core.rules
from core.dtos import AnalysisResultDto, FileInfoDto, SummaryDto, SlideResultDto, IssueDto
from core.enums import RuleId
from core.rules.base_rule import BaseRule
from core.utils.fonts_utils import get_used_base_fonts_in_pptx, check_font_substitution, get_used_fonts_from_pdf_path
from core.utils.main_utils import convert_ppt_to_pptx, convert_pptx_to_pdf

logger = logging.getLogger(__name__)


class AnalyserService:
    """
    Manages the analysis of uploaded files using various rules and handlers (for different format).
    """

    def __init__(self):
        self.handlers = {
            ".pptx": _analyse_pptx,
            ".ppt": _analyse_ppt,
            # Example for possible PDF analysis.
            ".pdf": _analyse_pdf,
        }
        self.all_rules = [rule_class() for rule_class in _REGISTERED_RULES]
        self.rules_map = {rule.rule_id.value: rule for rule in self.all_rules}

    def analyse(self, uploaded_file, rules_config: Optional[list[str]] = None) -> AnalysisResultDto:
        _, file_extension = os.path.splitext(uploaded_file.name)
        file_extension = file_extension.lower()

        handler_method = self.handlers.get(file_extension)
        if not handler_method:
            raise ValueError(f"Unsupported file type: '{file_extension}'")

        rules_to_apply = self._prepare_rules(rules_config)
        analysis_data = handler_method(uploaded_file, file_extension, rules_to_apply)

        return analysis_data

    def _prepare_rules(self, rules_config: list[str] | None) -> list[BaseRule]:
        """
        Parses the rules' configurations from the request and prepares the corresponding rule instances with parameters.
        """
        # Fallback to all default rules if no specific configuration is provided.
        if not rules_config:
            return self.all_rules

        prepared_rules: list[BaseRule] = []

        # Parse the config string to separate the rule ID from its optional parameter.
        for config_str in rules_config:
            rule_id, sep, param_value = config_str.partition(":")
            if not sep:
                param_value = None

            # Look up the base rule object. Skip silently (with debug log) if it doesn't exist.
            if (rule_template := self.rules_map.get(rule_id)) is None:
                logger.debug(f"Unknown rule id '{rule_id}' requested. Skipping.")
                continue

            # If no parameter provided or the rule itself doesn't accept any parameters; use the default rule template.
            if param_value is None or not (params := rule_template.get_parameters()):
                prepared_rules.append(rule_template)
                continue

            # Extract the expected name and type of the rule's first parameter.
            param_name = next(iter(params))
            param_type_name = params[param_name].get("type", "str")

            try:
                # Cast the string parameter value from the config to the required Python type.
                match param_type_name:
                    case "int":
                        typed_value = int(param_value)
                    case "float":
                        typed_value = float(param_value)
                    case "bool":
                        typed_value = param_value.lower() in ("true", "1", "yes")
                    case _:
                        typed_value = str(param_value)

                # Create an instance of the rule class, passing the typed parameter.
                # noinspection PyArgumentList
                new_rule_instance = type(rule_template)(**{param_name: typed_value})
                prepared_rules.append(new_rule_instance)

            except (ValueError, TypeError) as e:
                # If type casting fails; fall back to the default rule template.
                logger.warning(
                    f"Invalid parameter '{param_value}' for rule '{rule_id}'. "
                    f"Expected type {param_type_name}. Falling back to default. Error: {e}"
                )
                prepared_rules.append(rule_template)

        return prepared_rules


class RulesService:
    """
    Contains services to receive and manage a collection of rules.
    """

    def __init__(self):
        self.rules = [rule_class() for rule_class in _REGISTERED_RULES]

    def get_rules(self):
        return [
            {
                "id": rule.rule_id.value,
                "description": rule.rule_id.description,
                "parameters": rule.get_parameters(),
            }
            for rule in self.rules
        ]

    def get_rule_by_id(self, rule_id: str):
        for rule in self.rules:
            if rule.rule_id.value == rule_id:
                return {
                    "id": rule.rule_id.value,
                    "description": rule.rule_id.description,
                    "parameters": rule.get_parameters(),
                }
        return None


def _load_rules() -> list[type[BaseRule]]:
    """
    Dynamically loads all rule classes from the core.rules package.
    """
    loaded_classes: list[type[BaseRule]] = []
    # Iterate all files in rules package to find all rules.
    for module_info in pkgutil.walk_packages(core.rules.__path__, core.rules.__name__ + "."):
        try:
            module = importlib.import_module(module_info.name)

            # Get all rule classes.
            for _, cls in inspect.getmembers(module, inspect.isclass):
                # Get all subclasses of BaseRule.
                if issubclass(cls, BaseRule) and cls is not BaseRule:
                    # Check, if class is defined in the current module.
                    if cls.__module__ == module.__name__:
                        loaded_classes.append(cls)

        except ImportError as e:
            logger.exception(f"Error importing {module_info.name}: {e}")
        except Exception as e:
            logger.exception(f"Unexpected error loading rules from {module_info.name}: {e}")

    return loaded_classes


# Load rules once when the module is imported, so they are available for all instances of the services without needing to reload them.
_REGISTERED_RULES: list[type[BaseRule]] = _load_rules()


def _analyse_pptx(presentation_file, file_extension, rules: list[BaseRule]):
    """
    Analyses a PowerPoint (.pptx) file based on specified rules and generates a detailed analysis report including global and slide-specific issues.
    """
    file_name = Path(presentation_file.name).stem
    unique_id = uuid.uuid4().hex[:4]
    subfolder_name = f"{file_name[:10]}_{unique_id}"

    run_tmp_dir = Path(settings.TMP_DIR) / subfolder_name
    run_tmp_dir.mkdir(parents=True, exist_ok=True)

    # Create a temporary file to store the uploaded presentation.
    with NamedTemporaryFile(delete=False, suffix=file_extension, dir=run_tmp_dir) as tmp:
        if isinstance(presentation_file, Path):
            tmp.write(presentation_file.read_bytes())
        else:
            tmp.write(presentation_file.read())
        tmp_path = Path(tmp.name)
        logger.debug(f"Saved uploaded file with name {presentation_file.name} to {tmp_path}")

    pdf_path = convert_pptx_to_pdf(tmp_path, run_tmp_dir)
    try:
        pptx = Presentation(str(tmp_path))
        pdf = pdfium.PdfDocument(pdf_path)

        global_issues_collection: dict[RuleId, list[IssueDto]] = {}
        slide_issues_collection: dict[RuleId, dict[int, list[IssueDto]]] = {}

        # Check installed fonts after pdf conversion.
        pptx_fonts = get_used_base_fonts_in_pptx(pptx)
        logger.debug(f"Fonts used or embedded in PPTX: {pptx_fonts}")
        pdf_fonts = get_used_fonts_from_pdf_path(pdf_path)
        logger.debug(f"Fonts used or embedded in PDF: {pdf_fonts}")

        missing_fonts = check_font_substitution(pptx_fonts, pdf_fonts)
        if missing_fonts:
            logger.warning("=" * 60)
            logger.warning("FONT SUBSTITUTION DETECTED!")
            logger.warning("The following fonts are either missing from the Unoserver container or are declared as main fonts in the PPTX but aren't used in the slides:")
            for f in missing_fonts:
                logger.warning(f"- {f}")
            logger.warning("As a result, some applied analysis rules (contrast ratio, font size, etc.) may report false positives or false negatives.")
            logger.warning("Solution, if there are fonts missing in Docker setup: Mount the directory containing the missing fonts (.ttf/.otf files) into the unoserver container via your docker-compose.yml.")
            logger.warning("=" * 60)

        # Iterate through each rule and apply it to the presentation.
        for rule in rules:
            logger.debug(f"Applying rule {rule.rule_id.value} to presentation {presentation_file.name}")
            debug_start_rule = perf_counter()  # performance measurements for debugging
            result = rule.apply(pptx, pdf)
            logger.debug(f"Rule {rule.rule_id.value} applied in {(perf_counter() - debug_start_rule):.3f} s")

            global_issues_collection[rule.rule_id] = result.global_issues
            slide_issues_collection[rule.rule_id] = result.slide_issues

        file_info_dto = FileInfoDto(
            file_name=os.path.basename(presentation_file.name),
            file_size=f"{os.path.getsize(tmp_path) / 1024:.2f} KB",
            total_slides=len(pptx.slides),
        )

        total_slides_issues = sum(
            sum(len(issues) for issues in slide_issues.values()) for slide_issues in slide_issues_collection.values())
        total_sum_issues = sum(len(issues) for issues in global_issues_collection.values()) + total_slides_issues
        slides_with_issues = len({slide_number for issues in slide_issues_collection.values() for slide_number in issues
                                  if
                                  issues[slide_number]})
        summary_dto = SummaryDto(
            total_issues_found=total_sum_issues,
            slides_with_issues=slides_with_issues,
            rules_checked=[rule.rule_id.value for rule in rules],
        )

        global_issues_dtos: list[IssueDto] = []
        for rule_id, issues in global_issues_collection.items():
            for issue in issues:
                issue.rule_id = rule_id.value
                global_issues_dtos.append(issue)

        slide_results_dtos: list[SlideResultDto] = []
        for slide_number in range(1, len(pptx.slides) + 1):
            slide_issues = []
            for rule_id, issues in slide_issues_collection.items():
                if slide_number in issues:
                    slide_issues.extend(issues[slide_number])

            has_issues = len(slide_issues) > 0
            slide_result_dto = SlideResultDto(
                slide_number=slide_number,
                has_issues=has_issues,
                issues=slide_issues,
            )
            slide_results_dtos.append(slide_result_dto)

        analysis_result_dto = AnalysisResultDto(
            analysis_id=str(uuid.uuid4()),
            analysis_timestamp=datetime.now().isoformat(),
            file_info=file_info_dto,
            summary=summary_dto,
            global_issues=global_issues_dtos,
            slide_results=slide_results_dtos,
        )

        pdf.close()

        return analysis_result_dto

    finally:
        if settings.DEBUG:
            logger.debug(f"DEBUG mode is on. Not removing temporary files: {tmp_path} and {pdf_path}.")
        else:
            if tmp_path.exists():
                os.remove(tmp_path)
            if pdf_path.exists():
                os.remove(pdf_path)


def _analyse_ppt(presentation_file, file_extension, rules: list[BaseRule]):
    """
    Older .ppt files can be converted to .pptx first. Afterwards, the same logic as for .pptx is applied.
    """
    file_name = Path(presentation_file.name).stem
    unique_id = uuid.uuid4().hex[:4]
    subfolder_name = f"{file_name[:10]}_{unique_id}"

    run_tmp_dir = Path(settings.TMP_DIR) / subfolder_name
    run_tmp_dir.mkdir(parents=True, exist_ok=True)

    with NamedTemporaryFile(delete=False, suffix=file_extension, dir=run_tmp_dir) as tmp:
        if isinstance(presentation_file, Path):
            tmp.write(presentation_file.read_bytes())
        else:
            tmp.write(presentation_file.read())
        tmp_path = Path(tmp.name)
        logger.debug(f"Saved uploaded PPT file to temporary path: {tmp_path}")

    pptx_file = None
    try:
        pptx_path = convert_ppt_to_pptx(tmp_path, run_tmp_dir)
        pptx_file = Path(pptx_path)

        # Delegate to _analyse_pptx for the actual analysis.
        return _analyse_pptx(pptx_file, ".pptx", rules)

    finally:
        if settings.DEBUG:
            logger.debug(f"DEBUG mode is on. Not removing temporary files: {tmp_path} and {pptx_file}.")
        else:
            if tmp_path.exists():
                os.remove(tmp_path)
            if pptx_file and pptx_file.exists():
                os.remove(pptx_file)


# PDF Analysis as an example of how further document formats could be implemented.
def _analyse_pdf(presentation_file, file_extension, rules: list[BaseRule]):
    with NamedTemporaryFile(delete=False, suffix=file_extension, dir=settings.TMP_DIR) as tmp:
        tmp.write(presentation_file.read_bytes())
        tmp_path = Path(tmp.name)
        logger.debug(f"Saved uploaded PDF file to temporary path: {tmp_path}")

    try:
        with pdfium.PdfDocument(tmp_path) as pdf:
            global_issues_collection: dict[RuleId, list[IssueDto]] = {}
            slide_issues_collection: dict[RuleId, dict[int, list[IssueDto]]] = {}  # PDF pages act as slides.

            for rule in rules:
                result = rule.apply(None, pdf)
                global_issues_collection[rule.rule_id] = result.global_issues
                slide_issues_collection[rule.rule_id] = result.slide_issues

            file_info_dto = FileInfoDto(
                file_name=os.path.basename(presentation_file.name),
                file_size=os.path.getsize(tmp_path),
                total_slides=len(pdf.pages),
            )

            total_slides_issues = sum(sum(len(issues) for issues in slide_issues.values()) for slide_issues in
                                      slide_issues_collection.values())
            total_sum_issues = sum(len(issues) for issues in global_issues_collection.values()) + total_slides_issues
            slides_with_issues = len({slide_number for issues in slide_issues_collection.values() for slide_number in
                                      issues if issues[slide_number]}
                                     )

            summary_dto = SummaryDto(
                total_issues_found=total_sum_issues,
                slides_with_issues=slides_with_issues,
                rules_checked=[rule.rule_id.value for rule in rules],
            )

            global_issues_dtos: list[IssueDto] = []
            for rule_id, issues in global_issues_collection.items():
                for issue in issues:
                    issue.rule_id = rule_id.value
                    global_issues_dtos.append(issue)

            slide_results_dtos: list[SlideResultDto] = []
            for slide_number in range(1, len(pdf.pages) + 1):
                slide_issues = []
                for rule_id, issues in slide_issues_collection.items():
                    if slide_number in issues:
                        slide_issues.extend(issues[slide_number])

                has_issues = len(slide_issues) > 0
                slide_result_dto = SlideResultDto(
                    slide_number=slide_number,
                    has_issues=has_issues,
                    issues=slide_issues,
                )
                slide_results_dtos.append(slide_result_dto)

            analysis_result_dto = AnalysisResultDto(
                analysis_id=str(uuid.uuid4()),
                analysis_timestamp=datetime.now().isoformat(),
                file_info=file_info_dto,
                summary=summary_dto,
                global_issues=global_issues_dtos,
                slide_results=slide_results_dtos,
            )

            return analysis_result_dto
    finally:
        if tmp_path.exists():
            os.remove(tmp_path)
