from django.utils.translation import gettext as _
from pypdfium2 import PdfDocument
from pptx.presentation import Presentation

from core.dtos import IssueDto, RuleResultDto
from core.enums import RuleId
from core.rules.base_rule import BaseRule


class MaxImagesPerSlideRule(BaseRule):
    """
    Checks if the number of images on any slide exceeds a given maximum.
    """
    RULE_ID = RuleId.MEDIA_MAX_IMAGES_PER_SLIDE

    def __init__(self, max_images: int = 2):
        self.max_images = max_images

    def apply(self, pptx: Presentation | None, pdf: PdfDocument | None) -> RuleResultDto:
        global_issues: list[IssueDto] = []
        slide_issues: dict[int, list[IssueDto]] = {}

        if pptx:
            for i, slide in enumerate(pptx.slides, start=1):
                image_count = 0
                for shape in slide.shapes:
                    # Use this instead of shape.shape_type, because placeholders can also have images.
                    if hasattr(shape, "image") and shape.image is not None:
                        image_count += 1
                if image_count > self.max_images:
                    issue = IssueDto(
                        rule_id=self.RULE_ID.value,
                        message=_("Slide contains %(image_count)d images, which exceeds the recommended maximum of %(max)d.") % {
                            "image_count": image_count,
                            "max": self.max_images
                        },
                        details={
                            "image_count": image_count,
                            "max_images": self.max_images
                        }
                    )
                    slide_issues[i] = [issue]

        return RuleResultDto(
            global_issues=global_issues,
            slide_issues=slide_issues
        )
