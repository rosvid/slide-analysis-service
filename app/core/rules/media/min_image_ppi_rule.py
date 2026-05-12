from django.utils.translation import gettext as _
from pypdfium2 import PdfDocument
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.presentation import Presentation

from core.dtos import IssueDto, RuleResultDto
from core.enums import RuleId
from core.rules.base_rule import BaseRule


class MinImagePpiRule(BaseRule):
    """
    Checks if images on slides meet a minimum PPI (pixels per inch) requirement.
    """
    RULE_ID = RuleId.MEDIA_MIN_IMAGE_PPI

    def __init__(self, min_ppi: int = 150):
        self.min_ppi = min_ppi

    def apply(self, pptx: Presentation | None, pdf: PdfDocument | None) -> RuleResultDto:
        global_issues: list[IssueDto] = []
        slide_issues: dict[int, list[IssueDto]] = {}

        if pptx:
            for i, slide in enumerate(pptx.slides, start=1):
                for shape in slide.shapes:
                    # Check if shape has an image.
                    if hasattr(shape, "image") and shape.image is not None:
                        image = shape.image
                        image_pixels_w, image_pixels_h = image.size
                        shape_cm_w = shape.width.cm
                        shape_cm_h = shape.height.cm

                        # Extract the description (alternative text) from the XML.
                        descr_matches = shape.element.xpath(".//p:cNvPr/@descr")
                        description = descr_matches[0] if descr_matches else ""

                        # Calculate effective PPI (image could be stretched or shrunk).
                        effective_ppi_w = image_pixels_w / shape_cm_w
                        effective_ppi_h = image_pixels_h / shape_cm_h

                        # Check if either horizontal or vertical PPI is below the minimum.
                        if effective_ppi_w < self.min_ppi or effective_ppi_h < self.min_ppi:

                            # Grab position data for better identification (fallback to 0 if None).
                            shape_left = shape.left.cm if shape.left else 0.0
                            shape_top = shape.top.cm if shape.top else 0.0

                            # Use the lowest PPI for the main error message to keep it simple for the user.
                            lowest_ppi = int(min(effective_ppi_w, effective_ppi_h))

                            issue = IssueDto(
                                rule_id=self.RULE_ID.value,
                                message=_("Image '%(shape_name)s' has an effective resolution of %(lowest_ppi)s PPI, which is below the recommended minimum of %(min_ppi)s PPI.") % {
                                    "shape_name": shape.name,
                                    "lowest_ppi": lowest_ppi,
                                    "min_ppi": self.min_ppi
                                },
                                details={
                                    "shape_name": shape.name,
                                    "description": description,
                                    "effective_ppi_w": int(effective_ppi_w),
                                    "effective_ppi_h": int(effective_ppi_h),
                                    "min_ppi": self.min_ppi,
                                    "original_pixels": f"{image_pixels_w}x{image_pixels_h}",
                                    "rendered_cm": f"{round(shape_cm_w, 2)}x{round(shape_cm_h, 2)}",
                                    "position_cm": f"Left: {round(shape_left, 2)}, Top: {round(shape_top, 2)}"
                                }
                            )

                            if i not in slide_issues:
                                slide_issues[i] = []
                            slide_issues[i].append(issue)

        return RuleResultDto(
            global_issues=global_issues,
            slide_issues=slide_issues
        )
