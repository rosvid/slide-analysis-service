import logging

from django.utils.translation import gettext as _
from pypdfium2 import PdfDocument
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.presentation import Presentation

from core.dtos import IssueDto, RuleResultDto
from core.enums import RuleId
from core.rules.base_rule import BaseRule
from core.utils.pptx_utils import has_bullet
from core.utils.pptx_utils import has_master_bullets_at_level
from core.utils.pptx_utils import has_non_bullet
from core.utils.pptx_utils import has_placeholder_bullets_at_level
from core.utils.pptx_utils import has_placeholder_non_bullets_at_level

logger = logging.getLogger(__name__)


class MaxBulletPointsPerSlideRule(BaseRule):
    """
    Checks if the number of bullet points on any slide exceeds a given maximum.
    """
    RULE_ID = RuleId.TEXT_MAX_BULLET_POINTS_PER_SLIDE

    def __init__(self, max_bullet_points: int = 6):
        self.max_bullet_points = max_bullet_points

    def apply(self, pptx: Presentation | None, pdf: PdfDocument | None) -> RuleResultDto:
        global_issues: list[IssueDto] = []
        slide_issues: dict[int, list[IssueDto]] = {}

        if pptx:
            for i, slide in enumerate(pptx.slides, start=1):
                logger.debug(f"Analysing Slide {i} for bullet points.")
                bullet_points_count = 0
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue

                    slide_layout = slide.slide_layout
                    slide_master = slide_layout.slide_master

                    for paragraph in shape.text_frame.paragraphs:
                        # Skip empty paragraphs.
                        if not paragraph.text.strip():
                            continue

                        # Skip if paragraph is in a shape with placeholder which has no body type (skip titles, etc.).
                        if shape.is_placeholder:
                            placeholder_format = shape.placeholder_format
                            if placeholder_format.type not in [PP_PLACEHOLDER_TYPE.BODY, PP_PLACEHOLDER_TYPE.OBJECT]:  # body and obj type
                                continue

                        if has_non_bullet(paragraph._pPr):
                            continue

                        if has_bullet(paragraph._pPr):
                            logger.debug(f"Paragraph: '{paragraph.text}' has explicit bullet: {has_bullet(paragraph._pPr)}")
                            bullet_points_count += 1
                            continue

                        # If text is in a text box without bullets, skip layout and master checks.
                        # if shape.element.xpath(".//p:cNvSpPr[@txBox='1']"):
                        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                            continue

                        # Check if there is bullet style defined in slide layout for this level.
                        if shape.is_placeholder:
                            placeholder_format = shape.placeholder_format
                            ph_idx = placeholder_format.idx
                            placeholder = slide_layout.placeholders.get(ph_idx)
                            if placeholder and has_placeholder_non_bullets_at_level(placeholder, paragraph.level):
                                continue
                            if placeholder and has_placeholder_bullets_at_level(placeholder, paragraph.level):
                                logger.debug(f"Paragraph: '{paragraph.text}' has placeholder bullets at level {paragraph.level + 1}")
                                bullet_points_count += 1
                                continue

                        # Check for bullet style in slide master body style.
                        if has_master_bullets_at_level(slide_master, paragraph.level):
                            logger.debug(f"Paragraph: '{paragraph.text}' has master bullets at level {paragraph.level + 1}")
                            bullet_points_count += 1
                            continue

                logger.debug(f"Total bullet points counted on slide: {bullet_points_count}")

                if bullet_points_count > self.max_bullet_points:
                    issue = IssueDto(
                        rule_id=self.RULE_ID.value,
                        message=_("Slide contains %(count)d bullet points, which exceeds the recommended maximum of %(max)d.") % {
                            "count": bullet_points_count,
                            "max": self.max_bullet_points
                        }, details={
                            "bullet_points_count": bullet_points_count,
                            "max_bullet_points": self.max_bullet_points
                        }
                    )
                    slide_issues[i] = [issue]

        return RuleResultDto(
            global_issues=global_issues,
            slide_issues=slide_issues
        )
