from django.utils.translation import gettext as _
from pypdfium2 import PdfDocument
from pptx.presentation import Presentation

from core.dtos import IssueDto, RuleResultDto
from core.enums import RuleId
from core.rules.base_rule import BaseRule


class MaxWordsPerSlideRule(BaseRule):
    """
    Checks if the number of words on any slide exceeds a given maximum.
    """
    RULE_ID = RuleId.TEXT_MAX_WORDS_PER_SLIDE

    def __init__(self, max_words: int = 40):
        self.max_words = max_words

    def apply(self, pptx: Presentation | None, pdf: PdfDocument | None) -> RuleResultDto:
        global_issues: list[IssueDto] = []
        slide_issues: dict[int, list[IssueDto]] = {}

        if pptx:
            for i, slide in enumerate(pptx.slides, start=1):
                word_count = 0
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue

                    # Count words in the text frame.
                    text = shape.text_frame.text
                    word_count += len(text.split())

                if word_count > self.max_words:
                    issue = IssueDto(
                        rule_id=self.RULE_ID.value,
                        message=_("Slide contains %(word_count)d words, which exceeds the recommended maximum of %(max)d.") % {
                            "word_count": word_count,
                            "max": self.max_words
                        },
                        details={
                            "word_count": word_count,
                            "max_words": self.max_words
                        }
                    )
                    slide_issues[i] = [issue]

        return RuleResultDto(
            global_issues=global_issues,
            slide_issues=slide_issues
        )
