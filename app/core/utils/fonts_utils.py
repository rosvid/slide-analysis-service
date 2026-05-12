import logging
import re
from pathlib import Path

from pptx.opc.constants import RELATIONSHIP_TYPE
from pptx.oxml import parse_xml
from pptx.presentation import Presentation
from pypdf import PdfReader

logger = logging.getLogger(__name__)

IGNORED_FONTS = {
    "Symbol", "Wingdings", "Wingdings 2", "Wingdings 3", "Webdings",
    "Cambria Math", "Segoe UI Symbol", "Segoe UI Emoji", "Marlett"
}

# Regex: Looks for Light and Regular, which are sometimes used as suffixes in font names.
FONT_CLEANUP_PATTERN = re.compile(r"[\s\-]+(?:Light|Regular)\b", re.IGNORECASE)


def _camel_case_to_spaces(value: str) -> str:
    return re.compile(r"(?<=[a-z])(?=[A-Z])|(?<=[A-Z])(?=[A-Z][a-z])").sub(" ", value).strip()


def get_used_base_fonts_in_pptx(presentation: Presentation) -> set[str]:
    """
    Extracts all used base fonts from a PowerPoint presentation.
    """
    used_fonts = set()  # use a set to avoid duplicates
    needs_theme_fonts = False

    # fonts used in text frames
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Ignore empty runs without text.
                        if not run.text.strip():
                            continue

                        font_name = run.font.name
                        # Ignore fonts that are typically used for bullet points, symbols, emojis, icons.
                        if font_name and font_name not in IGNORED_FONTS:
                            used_fonts.add(FONT_CLEANUP_PATTERN.sub("", font_name).strip())
                        else:
                            needs_theme_fonts = True

    # embedded fonts
    font_list = presentation.element.xpath(".//p:embeddedFontLst/p:embeddedFont/p:font")
    for font_element in font_list:
        font_name = font_element.get("typeface")
        if font_name:
            used_fonts.add(FONT_CLEANUP_PATTERN.sub("", font_name).strip())

    if needs_theme_fonts:
        # Major and minor fonts from theme act as a fallback if no fonts are specified in text runs.
        theme_part = presentation.slide_master.part.part_related_by(RELATIONSHIP_TYPE.THEME)
        theme = parse_xml(theme_part.blob)

        major_font = theme.xpath(".//a:majorFont/a:latin/@typeface")
        minor_font = theme.xpath(".//a:minorFont/a:latin/@typeface")
        if major_font:
            used_fonts.add(FONT_CLEANUP_PATTERN.sub("", major_font[0]).strip())
        if minor_font:
            used_fonts.add(FONT_CLEANUP_PATTERN.sub("", minor_font[0]).strip())

    return used_fonts


def get_simple_font_name(font_string: str) -> str:
    """
    Simplifies a font name string by removing prefixes, suffixes and PostScript type ids.
    For example, "BAAAAA+TimesNewRomanPS-BoldMT" becomes "TimesNewRoman-Bold".
    """
    if "+" in font_string:
        font_string = font_string.split("+", 1)[1]

    family, style = font_string.split("-", 1) if "-" in font_string else (font_string, "")

    # Remove PostScript type identifiers only if they appear at the end.
    family = _camel_case_to_spaces(family.removesuffix("PSMT").removesuffix("PS"))
    style = _camel_case_to_spaces(style.removesuffix("MT").removesuffix("PS").removesuffix("Regular").removesuffix("Light"))

    return f"{family} {style}".strip()


def get_used_fonts_from_pdf_path(pdf_file_path: Path) -> set[str]:
    """
    Extracts a set of all unique fonts used within a specified PDF file.

    Note: `pypdf` package is used here because it's much faster than `pypdfium2` for this specific task.
    """
    fonts_used = set()
    reader = PdfReader(pdf_file_path)
    for i, page in enumerate(reader.pages):
        # verify whether the page contains resources and associated fonts
        if "/Resources" in page and "/Font" in page["/Resources"]:
            fonts_dict = page["/Resources"]["/Font"] # type: ignore

            # iterate through the identified font references
            for font_key in fonts_dict:
                font_obj = fonts_dict[font_key].get_object()

                # extract the font name from the BaseFont attribute
                if "/BaseFont" in font_obj:
                    raw_font_name = font_obj["/BaseFont"]
                    fonts_used.add(get_simple_font_name(raw_font_name))

    return fonts_used


def check_font_substitution(pptx_fonts: set[str], pdf_fonts: set[str]) -> set[str]:
    """
    Checks for potential font substitutions by comparing the fonts used in the PowerPoint presentation with
    those extracted from the PDF.

    If the fonts which are defined in the PPTX in the PDF are missing, it could mean that the font files are
    missing (e.g. on the server where the PDF was generated).
    """
    pdf_fonts_clean = {f.lower().replace(" ", "").replace("-", "") for f in pdf_fonts}
    missing_fonts = set()

    for pptx_font in pptx_fonts:
        clean_pptx_font = pptx_font.removesuffix("MT").removesuffix("PS")
        clean_pptx_font = clean_pptx_font.lower().replace(" ", "").replace("-", "")
        font_found = any(pdf_font.startswith(clean_pptx_font) for pdf_font in pdf_fonts_clean)

        if not font_found:
            missing_fonts.add(pptx_font)

    return missing_fonts
