import logging
import os
import uuid
from datetime import datetime
from pathlib import Path
from tempfile import NamedTemporaryFile

import pypdfium2 as pdfium
from django.conf import settings
from pptx import Presentation

from core.analysers.base_analyser import BaseAnalyser
from core.dtos import AnalysisResultDto, FileInfoDto, SummaryDto, SlideResultDto, IssueDto
from core.enums import RuleId
from core.rules.base_rule import BaseRule
from core.utils.fonts_utils import get_used_base_fonts_in_pptx, check_font_substitution, get_used_fonts_from_pdf_path
from core.utils.main_utils import convert_pptx_to_pdf

logger = logging.getLogger(__name__)


class PptxAnalyser(BaseAnalyser):
    def analyse(self, presentation_file, file_extension: str, rules: list[BaseRule]) -> AnalysisResultDto:
        file_name = Path(presentation_file.name).stem
        unique_id = uuid.uuid4().hex[:4]
        subfolder_name = f"{file_name[:10]}_{unique_id}"

        run_tmp_dir = Path(settings.TMP_DIR) / subfolder_name
        run_tmp_dir.mkdir(parents=True, exist_ok=True)

        # Create a temporary file to store the uploaded presentation.
        with NamedTemporaryFile(delete=False, suffix=file_extension, dir=run_tmp_dir) as tmp:
            if isinstance(presentation_file, Path):
                tmp.write(presentation_file.read_bytes())
            else:
                tmp.write(presentation_file.read())
            tmp_path = Path(tmp.name)
            logger.debug(f"Saved uploaded file with name {presentation_file.name} to {tmp_path}")

        pdf_path = convert_pptx_to_pdf(tmp_path, run_tmp_dir)
        try:
            pptx = Presentation(str(tmp_path))
            pdf = pdfium.PdfDocument(pdf_path)

            global_issues_collection: dict[RuleId, list[IssueDto]] = {}
            slide_issues_collection: dict[RuleId, dict[int, list[IssueDto]]] = {}

            # Check installed fonts after pdf conversion.
            pptx_fonts = get_used_base_fonts_in_pptx(pptx)
            logger.debug(f"Fonts used or embedded in PPTX: {pptx_fonts}")
            pdf_fonts = get_used_fonts_from_pdf_path(pdf_path)
            logger.debug(f"Fonts used or embedded in PDF: {pdf_fonts}")

            missing_fonts = check_font_substitution(pptx_fonts, pdf_fonts)
            if missing_fonts:
                logger.warning("=" * 60)
                logger.warning("FONT SUBSTITUTION DETECTED!")
                logger.warning("The following fonts are either missing from the Unoserver container or are declared as main fonts in the PPTX but aren't used in the slides:")
                for f in missing_fonts:
                    logger.warning(f"- {f}")
                logger.warning("As a result, some applied analysis rules (contrast ratio, font size, etc.) may report false positives or false negatives.")
                logger.warning("Solution, if there are fonts missing in Docker setup: Mount the directory containing the missing fonts (.ttf/.otf files) into the unoserver container via your docker-compose.yml.")
                logger.warning("=" * 60)

            # Iterate through each rule and apply it to the presentation.
            for rule in rules:
                logger.debug(f"Applying rule {rule.rule_id.value} to presentation {presentation_file.name}")
                # debug_start_rule = perf_counter()  # performance measurements for debugging
                result = rule.apply(pptx, pdf)
                # logger.debug(f"Rule {rule.rule_id.value} applied in {(perf_counter() - debug_start_rule):.3f} s")

                global_issues_collection[rule.rule_id] = result.global_issues
                slide_issues_collection[rule.rule_id] = result.slide_issues

            file_info_dto = FileInfoDto(
                file_name=os.path.basename(presentation_file.name),
                file_size=f"{os.path.getsize(tmp_path) / 1024:.2f} KB",
                total_slides=len(pptx.slides),
            )

            total_slides_issues = sum(
                sum(len(issues) for issues in slide_issues.values()) for slide_issues in
                slide_issues_collection.values())
            total_sum_issues = sum(len(issues) for issues in global_issues_collection.values()) + total_slides_issues
            slides_with_issues = len({slide_number for issues in slide_issues_collection.values() for slide_number in
                                      issues
                                      if
                                      issues[slide_number]})
            summary_dto = SummaryDto(
                total_issues_found=total_sum_issues,
                slides_with_issues=slides_with_issues,
                rules_checked=[rule.rule_id.value for rule in rules],
            )

            global_issues_dtos: list[IssueDto] = []
            for rule_id, issues in global_issues_collection.items():
                for issue in issues:
                    issue.rule_id = rule_id.value
                    global_issues_dtos.append(issue)

            slide_results_dtos: list[SlideResultDto] = []
            for slide_number in range(1, len(pptx.slides) + 1):
                slide_issues = []
                for rule_id, issues in slide_issues_collection.items():
                    if slide_number in issues:
                        slide_issues.extend(issues[slide_number])

                has_issues = len(slide_issues) > 0
                slide_result_dto = SlideResultDto(
                    slide_number=slide_number,
                    has_issues=has_issues,
                    issues=slide_issues,
                )
                slide_results_dtos.append(slide_result_dto)

            analysis_result_dto = AnalysisResultDto(
                analysis_id=str(uuid.uuid4()),
                analysis_timestamp=datetime.now().isoformat(),
                file_info=file_info_dto,
                summary=summary_dto,
                global_issues=global_issues_dtos,
                slide_results=slide_results_dtos,
            )

            pdf.close()

            return analysis_result_dto

        finally:
            if settings.DEBUG:
                logger.debug(f"DEBUG mode is on. Not removing temporary files: {tmp_path} and {pdf_path}.")
            else:
                if tmp_path.exists():
                    os.remove(tmp_path)
                if pdf_path.exists():
                    os.remove(pdf_path)
